#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation for CoupleSnap application
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "CoupleSnap"
subtitle.text = "Photo-First Messaging for Couples"

# Set title formatting
title_paragraph = title.text_frame.paragraphs[0]
title_paragraph.font.size = Pt(60)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 122, 255)

# Slide 2: Overview
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
title2 = slide2.shapes.title
content2 = slide2.placeholders[1]

title2.text = "What is CoupleSnap?"
content2.text = (
    "• Photo-first messaging app exclusively for couples\n"
    "• Every message must include a photo (primarily selfies)\n"
    "• Optional text overlay on photos\n"
    "• Dual-mode operation: Real-time & Asynchronous\n"
    "• Privacy-first: End-to-end encryption\n"
    "• No group features - just you and your partner"
)

# Slide 3: Key Features
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
title3 = slide3.shapes.title
content3 = slide3.placeholders[1]

title3.text = "Core Features"
content3.text = (
    "📷 Camera-First Interface\n"
    "   • Opens directly to camera view\n"
    "   • Front-facing camera by default\n"
    "   • Quick flip, flash, timer support\n\n"
    "💬 Photo Messaging\n"
    "   • Every message requires a photo\n"
    "   • Image compression & encryption\n"
    "   • Multiple resolution generation\n\n"
    "✍️ Text Overlay\n"
    "   • Draggable positioning\n"
    "   • Multiple fonts & styles\n"
    "   • Adjustable size and color\n\n"
    "🔔 Push Notifications\n"
    "   • Rich media previews\n"
    "   • Quick actions\n"
)

# Slide 4: Technology Stack
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
title4 = slide4.shapes.title
content4 = slide4.placeholders[1]

title4.text = "Technology Stack"
content4.text = (
    "Frontend:\n"
    "• React Native with TypeScript\n"
    "• Expo Camera\n"
    "• Zustand (State Management)\n"
    "• React Navigation\n\n"
    "Backend:\n"
    "• Firebase Authentication\n"
    "• Cloud Firestore (Real-time Database)\n"
    "• Firebase Storage (Media)\n"
    "• Firebase Cloud Functions\n"
    "• Firebase Cloud Messaging (Push)\n\n"
    "Security:\n"
    "• End-to-end encryption (AES-256)\n"
    "• Secure key storage"
)

# Slide 5: User Flow
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
title5 = slide5.shapes.title
content5 = slide5.placeholders[1]

title5.text = "User Experience Flow"
content5.text = (
    "1. Open App → Camera View\n"
    "2. Take Photo (Selfie)\n"
    "3. Add Optional Text Overlay\n"
    "   • Drag to position\n"
    "   • Choose font & style\n"
    "4. Send to Partner\n"
    "5. Partner Receives:\n"
    "   • Real-time if both online\n"
    "   • Push notification if offline\n"
    "6. View & React"
)

# Slide 6: Key Differentiators
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
title6 = slide6.shapes.title
content6 = slide6.placeholders[1]

title6.text = "What Makes Us Different"
content6.text = (
    "✅ Photo-Only Messaging\n"
    "   Every message requires a photo\n\n"
    "✅ Couple-Exclusive\n"
    "   Private space for two people only\n\n"
    "✅ Dual-Mode Operation\n"
    "   Smart switching between real-time & async\n\n"
    "✅ Privacy-First\n"
    "   End-to-end encryption, no public feeds\n\n"
    "✅ Intentional Communication\n"
    "   Photos create more meaningful exchanges"
)

# Slide 7: Architecture
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
title7 = slide7.shapes.title
content7 = slide7.placeholders[1]

title7.text = "System Architecture"
content7.text = (
    "Client Layer:\n"
    "• React Native App (iOS & Android)\n"
    "• 95% shared codebase\n"
    "• MVVM Pattern\n\n"
    "Backend Services:\n"
    "• Serverless Firebase Architecture\n"
    "• Auto-scaling & Multi-region\n"
    "• Real-time WebSocket connections\n\n"
    "Storage:\n"
    "• Firestore for data\n"
    "• Firebase Storage for media\n"
    "• CDN for fast delivery"
)

# Slide 8: Security & Privacy
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
title8 = slide8.shapes.title
content8 = slide8.placeholders[1]

title8.text = "Security & Privacy"
content8.text = (
    "🔒 End-to-End Encryption\n"
    "   • AES-256-GCM encryption\n"
    "   • RSA-2048 key pairs\n"
    "   • Forward secrecy\n\n"
    "🛡️ Authentication\n"
    "   • Phone number verification\n"
    "   • Biometric authentication support\n\n"
    "🔐 Data Protection\n"
    "   • Secure key storage (Keychain/Keystore)\n"
    "   • Certificate pinning\n"
    "   • Encrypted photo storage\n\n"
    "👥 Privacy Controls\n"
    "   • No public feeds\n"
    "   • No user discovery\n"
    "   • Read receipts & presence controls"
)

# Slide 9: Future Roadmap
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
title9 = slide9.shapes.title
content9 = slide9.placeholders[1]

title9.text = "Future Enhancements"
content9.text = (
    "Version 1.1:\n"
    "• Video messages (5-second clips)\n"
    "• Voice notes on photos\n"
    "• AR filters & effects\n"
    "• Scheduled messages\n\n"
    "Version 1.2:\n"
    "• Web app support\n"
    "• Apple Watch & Android Widgets\n"
    "• iPad optimized experience\n\n"
    "Version 2.0:\n"
    "• AI-powered features\n"
    "• Premium subscriptions\n"
    "• Physical photo products"
)

# Slide 10: Conclusion
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
title10 = slide10.shapes.title
content10 = slide10.placeholders[1]

title10.text = "Thank You"
content10.text = (
    "CoupleSnap\n"
    "Photo messaging for couples\n\n"
    "Making every moment visual\n"
    "and every exchange intentional\n\n"
    "Version 1.0"
)

# Save the presentation
output_file = "CoupleSnap_Presentation.pptx"
prs.save(output_file)
print(f"✅ Presentation created successfully: {output_file}")
print(f"📁 Location: {output_file}")


